"""
Generate slides/normalization_summary.pptx
Run:  python slides/build_normalization_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x2a, 0x3a)
BLUE   = RGBColor(0x3b, 0x7d, 0xd8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF7, 0xF9, 0xFC)
DGREY  = RGBColor(0x22, 0x22, 0x22)
YELLOW = RGBColor(0xFF, 0xFB, 0xE6)
AMBER  = RGBColor(0xE8, 0xC8, 0x40)
AMBER2 = RGBColor(0x5a, 0x40, 0x00)
CODEBG = RGBColor(0xFF, 0xF3, 0xC0)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_pt=0.75):
    shape = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, wrap=True):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tb.text_frame.word_wrap = wrap
    return tb.text_frame


def _run(para, text, size=11, bold=False, italic=False,
         color=DGREY, mono=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Courier New" if mono else "Arial"
    return run


def cell_runs(cell, parts, size=10):
    """
    parts: list of (text, bold, mono) tuples.
    Clears the cell and writes mixed-format runs.
    """
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_before = Pt(1)
    for text, bold, mono in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = DGREY
        r.font.name = "Courier New" if mono else "Arial"


# ─────────────────────────────────────────────────────────────────────────────
# Build slide
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)

PAD = Inches(0.35)

# ── Title bar ─────────────────────────────────────────────────────────────────
TITLE_H = Inches(0.52)
add_rect(sl, 0, 0, W, TITLE_H, fill=NAVY)
# Blue underline stripe
add_rect(sl, 0, TITLE_H, W, Inches(0.05), fill=BLUE)

tf_t = add_textbox(sl, PAD, Inches(0.08), W - PAD * 2, TITLE_H - Inches(0.08))
p = tf_t.paragraphs[0]
_run(p, "Image Normalization Strategies — SubCellAE Patch Prep",
     size=20, bold=True, color=WHITE)

# ── Rolling Ball banner ───────────────────────────────────────────────────────
BANNER_Y = TITLE_H + Inches(0.12)
BANNER_H = Inches(0.55)
add_rect(sl, PAD, BANNER_Y, W - PAD * 2, BANNER_H,
         fill=YELLOW, line_color=AMBER, line_pt=1.5)
add_rect(sl, PAD, BANNER_Y, Inches(0.06), BANNER_H, fill=AMBER)

tf_b = add_textbox(sl, PAD + Inches(0.15), BANNER_Y + Inches(0.07),
                   W - PAD * 2 - Inches(0.2), BANNER_H - Inches(0.1))
pb = tf_b.paragraphs[0]
pb.space_before = Pt(0)
_run(pb, "Pre-step (optional, all modes) — Rolling Ball Background Subtraction    ",
     size=11, bold=True, color=AMBER2)
_run(pb, "img = img − rolling_ball(img, radius=10)", size=10, mono=True, color=AMBER2)
_run(pb, "   ·   Removes slow background variation before normalization"
         "   ·   Percentile stats re-computed on RB-corrected pixels"
         "   ·   No zero-clipping applied",
     size=10, color=AMBER2)

# ── Table ─────────────────────────────────────────────────────────────────────
TBL_Y = BANNER_Y + BANNER_H + Inches(0.12)
TBL_H = H - TBL_Y - Inches(0.1)
TBL_W = W - PAD * 2

NCOLS = 5
NROWS = 10   # 1 header + 9 data rows

tbl = sl.shapes.add_table(NROWS, NCOLS,
                           int(PAD), int(TBL_Y),
                           int(TBL_W), int(TBL_H)).table

# Column widths: row-label col narrow, 4 mode cols equal
LABEL_W = TBL_W * 0.11
MODE_W  = (TBL_W - LABEL_W) / 4
for i, cw in enumerate([LABEL_W, MODE_W, MODE_W, MODE_W, MODE_W]):
    tbl.columns[i].width = int(cw)

# ── Header row ────────────────────────────────────────────────────────────────
headers = ["", "Dataset Percentile", "Image Percentile",
           "Cell Inside/Outside", "Cell Min-Max"]

for ci, hdr in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = hdr
    r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Arial"

# ── Data rows ─────────────────────────────────────────────────────────────────
# Each row: list of 5 cell specs, each spec = list of (text, bold, mono) tuples
# Use None for empty / simple string

def s(text, bold=False, mono=False):
    return (text, bold, mono)

rows = [
    # row label,  dataset,  image,  insideoutside,  minmax
    [
        [s("mode", bold=True)],
        [s('"dataset"', mono=True)],
        [s('"image"', mono=True)],
        [s('"cell_insideoutside"', mono=True)],
        [s('"cell_minmax"', mono=True)],
    ],
    [
        [s("Step 1", bold=True)],
        [s("p_lo, p_hi = percentile(\n  all pixels in set, 0.2/99.8)", mono=True)],
        [s("p_lo, p_hi = percentile(\n  this image, 1/99)", mono=True)],
        [s("int1 = img − median(outside)", mono=True)],
        [s("int1 = img − mean(outside)", mono=True)],
    ],
    [
        [s("Step 2", bold=True)],
        [s("out = clip(\n  (img−p_lo)/(p_hi−p_lo), 0,1)", mono=True)],
        [s("out = clip(\n  (img−p_lo)/(p_hi−p_lo), 0,1)", mono=True)],
        [s("int2 = int1 / mean(int1[inside])", mono=True)],
        [s("scale = percentile(\n  int1[inside], 99)", mono=True)],
    ],
    [
        [s("Step 3", bold=True)],
        [s("—")],
        [s("—")],
        [s("out = int2 / scale", mono=True), s("  (default 5)")],
        [s("out = clip(int1/scale, 0, 1)", mono=True)],
    ],
    [
        [s("Stats scope", bold=True)],
        [s("Whole dataset")],
        [s("Per image")],
        [s("Per image + mask")],
        [s("Per image + mask")],
    ],
    [
        [s("BG reference", bold=True)],
        [s("Global percentile")],
        [s("Per-image percentile")],
        [s("Median of outside pixels")],
        [s("Mean of outside pixels")],
    ],
    [
        [s("Output range", bold=True)],
        [s("[0, 1]", mono=True)],
        [s("[0, 1]", mono=True)],
        [s("~[0, 1]", mono=True), s(" (fixed scale)")],
        [s("[0, 1]", mono=True)],
    ],
    [
        [s("Mask needed", bold=True)],
        [s("No")],
        [s("No")],
        [s("Yes")],
        [s("Yes")],
    ],
    [
        [s("Notes", bold=True)],
        [s("Cross-image comparable. Requires full dataset pass. Recommended baseline.")],
        [s("No cross-image comparability. No pre-computation.")],
        [s("Robust to BG outliers via median. Scale keeps range predictable.")],
        [s("Matches Alana's notebook. 99th-pct avoids single-pixel compression.")],
    ],
]

for ri, row_data in enumerate(rows):
    bg = LGREY if ri % 2 == 1 else WHITE
    for ci, parts in enumerate(row_data):
        cell = tbl.cell(ri + 1, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.space_before = Pt(1)
        for text, bold, mono in parts:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(10)
            r.font.bold = bold
            r.font.color.rgb = DGREY
            r.font.name = "Courier New" if mono else "Arial"

# ─────────────────────────────────────────────────────────────────────────────
out = "slides/normalization_summary.pptx"
prs.save(out)
print(f"Saved → {out}")
