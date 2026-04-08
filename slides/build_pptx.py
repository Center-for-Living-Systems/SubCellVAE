"""
Generate slides/ML_meeting_presentation_subcellAE.pptx
Run:  python slides/build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x4a, 0x7a)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF2, 0xF6, 0xFB)
DGREY  = RGBColor(0x33, 0x33, 0x33)
YELLOW = RGBColor(0xFF, 0xF8, 0xE1)
AMBER  = RGBColor(0xF0, 0xA5, 0x00)
GREEN  = RGBColor(0x1a, 0x7a, 0x3a)
RED    = RGBColor(0x9a, 0x1a, 0x1a)
MUTED  = RGBColor(0x88, 0x88, 0x88)
BORD   = RGBColor(0xCC, 0xCC, 0xCC)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DGREY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_para(tf, text, size=16, bold=False, color=DGREY,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None):
    """Blue top bar with title (and optional subtitle)."""
    bar_h = Inches(1.1) if subtitle else Inches(0.85)
    add_rect(slide, 0, 0, W, bar_h, fill=NAVY)
    add_text(slide, title,
             Inches(0.4), Inches(0.08), W - Inches(0.8), Inches(0.6),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.62), W - Inches(0.8), Inches(0.4),
                 size=18, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF))
    return bar_h

def section_slide(title, subtitle=""):
    """Full-bleed navy section divider."""
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=NAVY)
    add_text(sl, title,
             Inches(1), Inches(2.5), W - Inches(2), Inches(1.2),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(sl, subtitle,
                 Inches(1), Inches(3.7), W - Inches(2), Inches(0.8),
                 size=24, color=RGBColor(0xCC, 0xDD, 0xFF),
                 align=PP_ALIGN.CENTER)
    return sl

def note_box(slide, text, l, t, w, h):
    """Yellow note box."""
    add_rect(slide, l, t, w, h, fill=YELLOW, line=AMBER)
    tb = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.08),
                                   w - Inches(0.24), h - Inches(0.16))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.size = Pt(13); run.font.color.rgb = RGBColor(0x5a, 0x40, 0x00)


def make_table(slide, headers, rows,
               l, t, w, h,
               col_widths=None,
               hdr_fill=NAVY, hdr_color=WHITE,
               row_size=13, hdr_size=13):
    """Add a formatted table."""
    ncols = len(headers)
    nrows = len(rows) + 1   # +1 for header
    tbl = slide.shapes.add_table(nrows, ncols, l, t, w, h).table

    # column widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(cw)
    # header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = hdr_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = hdr
        run.font.size = Pt(hdr_size); run.font.bold = True
        run.font.color.rgb = hdr_color
    # data rows
    for ri, row in enumerate(rows):
        bg = LGREY if ri % 2 == 1 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            # strip html span tags for plain text
            import re
            clean = re.sub(r'<[^>]+>', '', str(val))
            run.text = clean
            run.font.size = Pt(row_size)
            # colour coding
            if 'pending' in str(val) or val == '—':
                run.font.color.rgb = MUTED; run.font.italic = True
            elif any(c.isdigit() for c in str(val)):
                run.font.color.rgb = DGREY
            else:
                run.font.color.rgb = DGREY
    return tbl

def bullets(slide, items, l, t, w, h, size=16, title=None, title_size=17):
    """Text box with optional bold title and bullet list."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = title
        run.font.size = Pt(title_size); run.font.bold = True
        run.font.color.rgb = NAVY
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "  – ") + item.lstrip("  ")
        run.font.size = Pt(size); run.font.color.rgb = DGREY


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(3.2), W, Inches(1.2), fill=RGBColor(0x0e, 0x2a, 0x50))
add_text(sl, "SubCellAE — ML Meeting",
         Inches(1), Inches(1.8), W - Inches(2), Inches(1.0),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Where We Stand",
         Inches(1), Inches(2.85), W - Inches(2), Inches(0.7),
         size=28, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
add_text(sl, "Can we use autoencoders to predict FA type and subcellular position\n"
             "from fluorescence microscopy patches?",
         Inches(1.5), Inches(4.8), W - Inches(3), Inches(1.2),
         size=18, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER, italic=True)

# =============================================================================
# SLIDE 2 — What Data Are We Looking At?
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "What Data Are We Looking At?")
top = bh + Inches(0.2)
col_w = (W - Inches(0.8)) / 2
L1 = Inches(0.4); L2 = L1 + col_w + Inches(0.1)

bullets(sl, [
    "Control — unperturbed cells",
    "Y-comp — Y-27632 (Rho kinase inhibitor)",
], L1, top, col_w, Inches(0.8), title="Experimental conditions:")

# channel table
make_table(sl,
    ["Channel", "Marker"],
    [["Ch 1", "Paxillin (PAX) — primary FA marker"],
     ["Ch 2", "eGFP-Zyxin (Z)"],
     ["Ch 3", "Phalloidin / Actin (A)"],
     ["Ch 4", "Vinculin (V)"]],
    L1, top + Inches(1.1), col_w, Inches(1.6),
    col_widths=[Inches(1.1), col_w - Inches(1.1)])

add_text(sl, "Imaging: Multi-channel .czi fluorescence microscopy",
         L1, top + Inches(2.85), col_w, Inches(0.4), size=13, italic=True, color=MUTED)

# right column
bullets(sl, [
    "Nascent Adhesion", "Focal Complex", "Focal Adhesion",
    "Fibrillar Adhesion", "No Adhesion",
], L2, top, col_w, Inches(1.6), title="FA type (5 classes):", size=15)
bullets(sl, [
    "Cell Protruding Edge", "Cell Periphery/other",
    "Lamella", "Cell Body",
], L2, top + Inches(1.75), col_w, Inches(1.3), title="Position (4 classes):", size=15)
add_text(sl, "Scale: ~5,800 patches per condition\n~60 manually labelled patches (new validation set)",
         L2, top + Inches(3.2), col_w, Inches(0.8), size=14, color=DGREY)

# =============================================================================
# SLIDE 3 — Normalization
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Normalization", "Per-dataset percentile stretch")
top = bh + Inches(0.25)
L = Inches(0.5); TW = W - Inches(1.0)

bullets(sl, [
    "Compute 0.2nd–99.8th percentile intensity across all images in the same condition set",
    "Apply a single linear mapping:  I_norm = clip( (I − p_low) / (p_high − p_low),  0, 1 )",
    "Applied uniformly → all images share the same intensity scale",
], L, top, TW, Inches(1.4), size=16)

add_text(sl, "Why this approximates per-cell normalization:",
         L, top + Inches(1.6), TW, Inches(0.4), size=16, bold=True, color=NAVY)
bullets(sl, [
    "Cells tend to fill the field of view → dataset percentiles are driven by cell signal, not background",
    "Avoids patch-to-patch intensity variation that would confuse the AE",
], L, top + Inches(2.05), TW, Inches(0.9), size=16)

note_box(sl,
    "Optional pre-step (all modes): Rolling Ball Background Subtraction\n"
    "img = img − rolling_ball(img, radius=10)  ·  Removes slow background variation  ·  "
    "Percentile stats re-computed on RB-corrected pixels",
    L, top + Inches(3.1), TW, Inches(0.85))

add_text(sl, "Current scope: normalization applied to the paxillin channel only",
         L, top + Inches(4.1), TW, Inches(0.4), size=14, italic=True, color=MUTED)

# =============================================================================
# SLIDE 4 — Autoencoder Basic Rules
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Autoencoder — Basic Rules")
top = bh + Inches(0.2)
col_w = (W - Inches(0.9)) / 2
L1 = Inches(0.4); L2 = L1 + col_w + Inches(0.1)

bullets(sl, [
    "Size: 32 × 32 pixels",
    "Grid-based, inside cell mask (≥ 40% overlap)",
    "Padded image (64 px) → deterministic coordinates",
], L1, top, col_w, Inches(1.3), title="Patch extraction:", size=15)

bullets(sl, [
    "3× Conv encoder → 8-dim latent z",
    "3× ConvTranspose decoder → Hardtanh [0,1] output",
    "Single channel input (paxillin)",
], L1, top + Inches(1.5), col_w, Inches(1.2), title="Architecture:", size=15)

bullets(sl, [
    "Group-aware 80/20 split — all patches from same image stay together",
    "Adam optimizer, LR = 0.001, cosine decay",
    "Batch = 128, up to 500 epochs",
], L2, top, col_w, Inches(1.3), title="Training:", size=15)

bullets(sl, [
    "8 distance-to-cell-edge features (d00–d07, rotation-invariant)",
    "Used as optional supplement to latent features in downstream classifier",
], L2, top + Inches(1.5), col_w, Inches(1.0), title="Additional features:", size=15)

# =============================================================================
# SLIDE 5 — Research Question
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Research Question")
top = bh + Inches(0.3)
L = Inches(0.5); TW = W - Inches(1.0)

add_text(sl, "Can we use an autoencoder to predict FA type and subcellular position?",
         L, top, TW, Inches(0.7),
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "Approach: Train AE on all patches → extract 8-dim latent z → train LightGBM classifier on z",
         L, top + Inches(0.75), TW, Inches(0.45), size=16, color=DGREY)

make_table(sl,
    ["Strategy", "Idea"],
    [["Baseline AE",         "Pure reconstruction — latent organizes by image content"],
     ["Semi-supervised AE",  "Add classification head during AE training using labelled patches"],
     ["Contrastive AE",      "Pull augmented views of the same patch together (NT-Xent loss)"],
     ["Supervised Contrastive AE", "Like contrastive, but same-class patches are positive pairs"]],
    L, top + Inches(1.3), TW, Inches(2.0),
    col_widths=[Inches(3.2), TW - Inches(3.2)])

add_text(sl, "Feature sets:   lat8  (latent only)   ·   lat8 + dist8  (latent + 8 distance features)",
         L, top + Inches(3.5), TW, Inches(0.45), size=15, color=DGREY)

# =============================================================================
# SLIDE 6 — Train/Val Split
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Classifier Train/Val Split — Two Strategies")
top = bh + Inches(0.2)
col_w = (W - Inches(0.9)) / 2
L1 = Inches(0.4); L2 = L1 + col_w + Inches(0.1)

bullets(sl, [
    "The split column in latents.csv is set during AE training: all patches from the same image file go to the same split (80% train / 20% val)",
    "Classifier inherits this split exactly",
    "No image-level leakage: train and val patches come from different microscopy images",
], L1, top, col_w, Inches(2.5), title="from_csv (group-aware):", size=14)

bullets(sl, [
    "A fresh stratified random split is drawn on individual patches, ignoring image origin",
    "Patches from the same image can appear in both train and val",
], L2, top, col_w, Inches(1.5), title="stratified (random):", size=14)

bullets(sl, [
    "Patches from the same cell image share: global illumination, cell shape, correlated local textures",
    "With stratified split, the classifier can memorize image-level cues rather than FA morphology",
], L2, top + Inches(1.8), col_w, Inches(1.5), title="Information leakage risk:", size=14)

note_box(sl,
    "Consequence: stratified validation accuracy is optimistically biased — the model appears to "
    "perform well but may fail on new experiments (e.g. Dy2).  "
    "from_csv gives a more honest estimate of cross-image generalization.",
    Inches(0.4), top + Inches(3.5), W - Inches(0.8), Inches(0.85))

# =============================================================================
# SECTION — Results
# =============================================================================
section_slide("Results", "Using Paxillin Only")

# =============================================================================
# SLIDE 7 — Training Set Performance
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Paxillin Only — Training Set Performance  (lat8)",
                  "Accuracy on held-out validation split (same experiment as training)")
top = bh + Inches(0.25)
L = Inches(0.5); TW = W - Inches(1.0)

make_table(sl,
    ["Model", "Position", "FA Classification"],
    [["Baseline AE",              "64%", "62%"],
     ["SemiSup AE (FA only)",     "60%", "92%"],
     ["SemiSup AE (FA + Pos)",    "95%", "93%"],
     ["Contrastive AE (NT-Xent)", "58%", "41%"],
     ["SupCon AE (+ flips)",      "52%", "37%"],
     ["SupCon AE (noise only)",   "54%", "41%"]],
    L, top, TW, Inches(2.6),
    col_widths=[Inches(4.5), Inches(2.8), Inches(2.8)])

note_box(sl,
    "SemiSup models trained on this dataset — high accuracy expected. "
    "Contrastive models receive no label supervision → lower classification accuracy but potentially better generalization.",
    L, top + Inches(2.8), TW, Inches(0.75))

# =============================================================================
# SLIDE 8 — New Data Validation
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Paxillin Only — New Data Validation  (Dy2, lat8)",
                  "Apply frozen trained models to new control images (Dy2 dataset) — 60 manually labelled patches")
top = bh + Inches(0.25)
L = Inches(0.5); TW = W - Inches(1.0)

make_table(sl,
    ["Model", "Position (Dy2)", "FA Classification (Dy2)"],
    [["Baseline AE",              "47%",     "47%"],
     ["SemiSup AE (FA only)",     "35%",     "37%"],
     ["SemiSup AE (FA + Pos)",    "41%",     "34%"],
     ["Contrastive AE (NT-Xent)", "pending", "pending"],
     ["SupCon AE (+ flips)",      "pending", "pending"],
     ["SupCon AE (noise only)",   "pending", "pending"]],
    L, top, TW, Inches(2.6),
    col_widths=[Inches(4.5), Inches(2.8), Inches(2.8)])

note_box(sl,
    "Semi-supervised models show a large drop from training to new data — see notes slide.",
    L, top + Inches(2.8), TW, Inches(0.55))

# =============================================================================
# SLIDE 9 — Combined Summary
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Paxillin Only — Combined Summary")
top = bh + Inches(0.2)
L = Inches(0.5); TW = W - Inches(1.0)

make_table(sl,
    ["Model", "Pos (train)", "FA (train)", "Pos (Dy2)", "FA (Dy2)"],
    [["Baseline AE",      "64%", "62%", "47%", "47%"],
     ["SemiSup (FA)",     "60%", "92%", "35%", "37%"],
     ["SemiSup (Both)",   "95%", "93%", "41%", "34%"],
     ["Contrastive",      "58%", "41%", "—",   "—"],
     ["SupCon + flips",   "52%", "37%", "—",   "—"],
     ["SupCon no flip",   "54%", "41%", "—",   "—"]],
    L, top, TW, Inches(2.4),
    col_widths=[Inches(3.0), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)])

add_text(sl, "Adding distance features  (lat8 + dist8) — training set:",
         L, top + Inches(2.6), TW, Inches(0.4), size=15, bold=True, color=NAVY)

make_table(sl,
    ["Model", "Position", "FA Classification"],
    [["Baseline AE",    "76%", "67%"],
     ["SemiSup (Both)", "94%", "94%"]],
    L, top + Inches(3.1), Inches(7.5), Inches(0.95),
    col_widths=[Inches(3.5), Inches(2.0), Inches(2.0)])

add_text(sl, "Distance features improve Baseline position by +12 pp;  SemiSup unchanged.",
         L, top + Inches(4.2), TW, Inches(0.4), size=14, italic=True, color=DGREY)

# =============================================================================
# SLIDE 10 — Notes SemiSup
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Notes — Semi-Supervised AE")
top = bh + Inches(0.25)
L = Inches(0.5); TW = W - Inches(1.0)

bullets(sl, [
    "Two classification heads on top of shared encoder: FA type + subcellular position",
    "Classification loss active only on labelled patches (~20% of training set)",
    "Two-phase training: 200 epochs reconstruction-only warmup, then classification heads activated",
    "Loss:  λ_recon · MSE  +  λ_cls · CE(FA)  +  λ_cls_2 · CE(Position)",
], L, top, TW, Inches(1.8), title="Training setup:", size=15)

add_text(sl, "Why the large train → Dy2 gap?",
         L, top + Inches(2.0), TW, Inches(0.4), size=16, bold=True, color=NAVY)

note_box(sl,
    "The reference run (test_run_overfit_20260322) was intentionally set up to probe overfitting. "
    "The model has seen training images many times with label supervision — the latent space is "
    "shaped around those specific cells. On new images from a different imaging session, the class "
    "structure does not transfer cleanly.",
    L, top + Inches(2.5), TW, Inches(1.0))

bullets(sl, [
    "Test with stricter regularization / lower lambda_cls",
    "Compare against baseline on more new-data images (currently only 60 labelled patches)",
    "Run contrastive variants on Dy2 to see if unsupervised methods generalize better",
], L, top + Inches(3.7), TW, Inches(1.3), title="Next steps:", size=15)

# =============================================================================
# SECTION — Multi-Channel Results
# =============================================================================
section_slide("Multi-Channel Results", "(In Progress)")

# =============================================================================
# SLIDE 11 — P+A Training Set
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Paxillin + Actin (Ch1 + Ch3) — Training Set",
                  "16-dimensional latent space")
top = bh + Inches(0.2)
L = Inches(0.5); TW = W - Inches(1.0)
col_w = (TW - Inches(0.2)) / 2

add_text(sl, "Latent only (lat16):",
         L, top, col_w, Inches(0.4), size=14, bold=True, color=NAVY)
make_table(sl,
    ["Model", "Position", "FA Cls."],
    [["Baseline AE",          "58%", "44%"],
     ["SemiSup (FA only)",    "49%", "17%"],
     ["SemiSup (Pos only)",   "44%", "20%"],
     ["SemiSup (FA + Pos)",   "51%", "25%"]],
    L, top + Inches(0.45), col_w, Inches(1.7),
    col_widths=[Inches(3.0), Inches(1.5), Inches(1.5)])

R = L + col_w + Inches(0.2)
add_text(sl, "Latent + distance (lat16 + dist8):",
         R, top, col_w, Inches(0.4), size=14, bold=True, color=NAVY)
make_table(sl,
    ["Model", "Position", "FA Cls."],
    [["Baseline AE",          "72%", "44%"],
     ["SemiSup (FA only)",    "75%", "39%"],
     ["SemiSup (Pos only)",   "70%", "38%"],
     ["SemiSup (FA + Pos)",   "69%", "36%"]],
    R, top + Inches(0.45), col_w, Inches(1.7),
    col_widths=[Inches(3.0), Inches(1.5), Inches(1.5)])

note_box(sl,
    "Distance features give +14 pp position boost for Baseline (58→72%).  "
    "FA accuracy remains low (~44%) — actin channel does not add FA-type signal.  "
    "Dy2 new-data validation: pending.",
    L, top + Inches(2.4), TW, Inches(0.85))

# =============================================================================
# SLIDE 12 — P+Z  &  P+V  (combined pending slide)
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Paxillin + Zyxin  /  Paxillin + Vinculin — Pending")
top = bh + Inches(0.25)
L = Inches(0.5); TW = W - Inches(1.0)

add_text(sl, "Paxillin + Zyxin (Ch1 + Ch2):",
         L, top, TW, Inches(0.4), size=15, bold=True, color=NAVY)
make_table(sl,
    ["Model", "Position", "FA Cls.", "Pos (Dy2)", "FA (Dy2)"],
    [["Baseline AE",        "—", "—", "—", "—"],
     ["SemiSup (FA only)",  "—", "—", "—", "—"],
     ["SemiSup (FA + Pos)", "—", "—", "—", "—"],
     ["Contrastive AE",     "—", "—", "—", "—"]],
    L, top + Inches(0.45), TW, Inches(1.55),
    col_widths=[Inches(3.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)])

add_text(sl, "Paxillin + Vinculin (Ch1 + Ch4):",
         L, top + Inches(2.2), TW, Inches(0.4), size=15, bold=True, color=NAVY)
make_table(sl,
    ["Model", "Position", "FA Cls.", "Pos (Dy2)", "FA (Dy2)"],
    [["Baseline AE",        "—", "—", "—", "—"],
     ["SemiSup (FA only)",  "—", "—", "—", "—"],
     ["SemiSup (FA + Pos)", "—", "—", "—", "—"],
     ["Contrastive AE",     "—", "—", "—", "—"]],
    L, top + Inches(2.65), TW, Inches(1.55),
    col_widths=[Inches(3.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)])

# =============================================================================
# SLIDE 13 — Summary
# =============================================================================
sl = prs.slides.add_slide(BLANK)
bh = slide_header(sl, "Where We Stand — Summary")
top = bh + Inches(0.2)
L = Inches(0.5); TW = W - Inches(1.0)

make_table(sl,
    ["", "Done", "In progress"],
    [["Data",                  "Paxillin-only pipeline, patch extraction, normalization", "Multi-channel (P+Z, P+V) runs"],
     ["Models",                "Baseline, SemiSup (FA/Pos/Both), Contrastive, SupCon",   "Hyperparameter tuning; contrastive Dy2 eval"],
     ["Training eval (pax)",   "All 6 variants × lat8 + lat8dist8",                      "—"],
     ["Training eval (P+A)",   "Baseline + SemiSup × lat16 + lat16dist8",                "P+Z, P+V; contrastive multichannel"],
     ["New-data (Dy2)",        "Baseline + SemiSup (paxillin-only)",                     "All multichannel; more labelled patches"]],
    L, top, TW, Inches(2.3),
    col_widths=[Inches(2.8), Inches(5.5), Inches(4.0)], row_size=13)

bullets(sl, [
    "SemiSup AE achieves high training accuracy (92–95%) but does not generalize to Dy2 in the overfit test run",
    "Baseline AE (no label supervision) is more consistent across train and Dy2  (~47–64%)",
    "Distance-to-edge features consistently help position accuracy  (+12–14 pp across all models)",
    "P+A multichannel: position similar to paxillin-only baseline, FA accuracy lower (44% vs 62%) — actin doesn't add FA-type signal",
], L, top + Inches(2.5), TW, Inches(2.0), title="Key findings so far:", size=14)

note_box(sl,
    "Open question: Is the train→Dy2 drop a fundamental limitation of semi-supervision, "
    "or can we fix it with better regularization / more labelled data?",
    L, top + Inches(4.65), TW, Inches(0.65))

# =============================================================================
# SLIDE 14 — Thank you
# =============================================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_text(sl, "Thank you",
         Inches(1), Inches(2.8), W - Inches(2), Inches(1.0),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Questions?",
         Inches(1), Inches(3.9), W - Inches(2), Inches(0.6),
         size=24, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# =============================================================================
# Save
# =============================================================================
out = "slides/ML_meeting_presentation_subcellAE.pptx"
prs.save(out)
print(f"Saved → {out}")
