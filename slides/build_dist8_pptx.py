"""
Generate slides/dist8_feature_description.pptx
Run:  python slides/build_dist8_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x4a, 0x7a)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF2, 0xF6, 0xFB)
DGREY  = RGBColor(0x33, 0x33, 0x33)
YELLOW = RGBColor(0xFF, 0xF8, 0xE1)
AMBER  = RGBColor(0xF0, 0xA5, 0x00)
MONO   = RGBColor(0x22, 0x22, 0x22)
CODEBG = RGBColor(0xF4, 0xF4, 0xF4)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_pt=0.75):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = _Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def _run(para, text, size=18, bold=False, italic=False, color=DGREY,
         font_name="Arial", mono=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Courier New" if mono else font_name
    return run


def add_textbox(slide, l, t, w, h, wrap=True):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tb.text_frame.word_wrap = wrap
    return tb.text_frame


def make_table(slide, headers, rows, l, t, w, h,
               col_widths=None, hdr_size=13, row_size=13):
    ncols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, ncols, int(l), int(t), int(w), int(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(cw)
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.size = Pt(hdr_size); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Arial"
    for ri, row in enumerate(rows):
        bg = LGREY if ri % 2 == 1 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.size = Pt(row_size)
            r.font.color.rgb = DGREY
            r.font.name = "Courier New" if ci == 0 else "Arial"


# ─────────────────────────────────────────────────────────────────────────────
# Single slide
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)

PAD   = Inches(0.45)
TITLE_H = Inches(0.55)

# ── Title bar ────────────────────────────────────────────────────────────────
add_rect(sl, 0, 0, W, TITLE_H + Inches(0.1), fill=NAVY)
tb = add_textbox(sl, PAD, Inches(0.06), W - PAD * 2, TITLE_H)
p = tb.paragraphs[0]
_run(p, "Distance-to-Edge Feature Vector (dist8)",
     size=28, bold=True, color=WHITE)

TOP  = TITLE_H + Inches(0.25)
BOT  = H - Inches(0.2)
BODY = BOT - TOP

# Two columns: left 55%, right 45%
GAP   = Inches(0.3)
L_W   = Inches(7.0)
R_X   = PAD + L_W + GAP
R_W   = W - R_X - PAD

# ── LEFT COLUMN ──────────────────────────────────────────────────────────────
tf = add_textbox(sl, PAD, TOP, L_W, BODY)

def para(tf, space_before=6):
    """Add a new paragraph to tf, return it."""
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.space_before = _Pt(space_before)
    return p

p0 = tf.paragraphs[0]
p0.space_before = Pt(0)
_run(p0, "What it encodes", size=17, bold=True, color=NAVY)

p1 = para(tf, 4)
_run(p1,
     "Each patch center casts ", size=16, color=DGREY)
_run(p1, "8 rays", size=16, bold=True, color=DGREY)
_run(p1,
     " at evenly-spaced angles (0°, 45°, 90°, …, 315°) outward through the "
     "cell mask, recording the pixel distance to the nearest cell boundary in "
     "each direction.",
     size=16, color=DGREY)

p2 = para(tf, 10)
_run(p2, "Rotation invariance", size=17, bold=True, color=NAVY)

p3 = para(tf, 4)
_run(p3, "The 8 distances are ", size=16, color=DGREY)
_run(p3, "cyclically shifted", size=16, bold=True, color=DGREY)
_run(p3,
     " so that the direction of minimum distance is always placed first "
     "(d00). This makes the feature independent of cell orientation.",
     size=16, color=DGREY)

p4 = para(tf, 10)
_run(p4, "Result", size=17, bold=True, color=NAVY)

p5 = para(tf, 4)
_run(p5,
     "An 8-dimensional vector [d00, d01, …, d07] where:",
     size=16, color=DGREY)

for bullet in [
    "d00 = shortest ray (closest edge)",
    "d04 ≈ ray roughly opposite to d00",
]:
    pb = para(tf, 2)
    pb.level = 1
    _run(pb, "• " + bullet, size=15, color=DGREY)

p6 = para(tf, 10)
_run(p6, "Used as", size=17, bold=True, color=NAVY)

p7 = para(tf, 4)
_run(p7,
     "Supplement to the AE latent z for LightGBM classification — "
     "combined as ", size=16, color=DGREY)
_run(p7, "lat8 + dist8", size=16, bold=True, mono=True, color=DGREY)
_run(p7, "  or  ", size=16, color=DGREY)
_run(p7, "lat16 + dist8", size=16, bold=True, mono=True, color=DGREY)

# ── RIGHT COLUMN ─────────────────────────────────────────────────────────────

# ASCII diagram box
DIAG_H = Inches(2.0)
add_rect(sl, R_X, TOP, R_W, DIAG_H, fill=CODEBG, line_color=RGBColor(0xCC, 0xCC, 0xCC))

diag_lines = [
    "         d02  d01",
    "          |  /",
    "    d03 --•-- d00  ← min (always d00)",
    "          |  \\",
    "         d04  d07",
    "           ↓",
    "     cell boundary",
]

tf_d = add_textbox(sl, R_X + Inches(0.15), TOP + Inches(0.12),
                   R_W - Inches(0.3), DIAG_H - Inches(0.2))
first = True
for line in diag_lines:
    p = tf_d.paragraphs[0] if first else tf_d.add_paragraph()
    first = False
    p.space_before = Pt(0)
    _run(p, line, size=13, mono=True, color=MONO)

# Feature table
TBL_Y = TOP + DIAG_H + Inches(0.2)
TBL_H = Inches(1.7)
make_table(
    sl,
    ["Feature", "Meaning"],
    [
        ("d00", "distance to nearest edge"),
        ("d01–d03", "distances at +45°, +90°, +135°"),
        ("d04", "distance to opposite side"),
        ("d05–d07", "remaining directions"),
    ],
    R_X, TBL_Y, R_W, TBL_H,
    col_widths=[R_W * 0.32, R_W * 0.68],
    hdr_size=13, row_size=13,
)

# Note box
NOTE_Y = TBL_Y + TBL_H + Inches(0.2)
NOTE_H = H - NOTE_Y - Inches(0.15)
add_rect(sl, R_X, NOTE_Y, R_W, NOTE_H, fill=YELLOW, line_color=AMBER)
# Amber left bar
add_rect(sl, R_X, NOTE_Y, Inches(0.07), NOTE_H, fill=AMBER)

tf_n = add_textbox(sl, R_X + Inches(0.18), NOTE_Y + Inches(0.1),
                   R_W - Inches(0.25), NOTE_H - Inches(0.15))
pn = tf_n.paragraphs[0]
pn.space_before = Pt(0)
_run(pn, "Why it helps:  ", size=13, bold=True, color=DGREY)
_run(pn,
     "Subcellular position (edge vs. lamella vs. cell body) is strongly "
     "encoded in the relative distances — a patch at the periphery has one "
     "very short ray; a central patch has long rays in all directions.",
     size=13, color=DGREY)

# ─────────────────────────────────────────────────────────────────────────────
out = "slides/dist8_feature_description.pptx"
prs.save(out)
print(f"Saved → {out}")
